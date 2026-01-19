from pptx import Presentation
from pptx.util import Inches
from image_engine import get_image
import requests
from io import BytesIO

def create_ppt(data):
    prs = Presentation()

    # TITLE
    t = prs.slides.add_slide(prs.slide_layouts[0])
    t.shapes.title.text = data["title"]

    for s in data["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s["title"]

        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        for b in s["bullets"]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 1

        img_url = get_image(s["image_query"])
        if img_url:
            img = requests.get(img_url).content
            slide.shapes.add_picture(
                BytesIO(img),
                Inches(5.5), Inches(1.5),
                width=Inches(4)
            )

    prs.save("ai_slides.pptx")
    return "ai_slides.pptx"
